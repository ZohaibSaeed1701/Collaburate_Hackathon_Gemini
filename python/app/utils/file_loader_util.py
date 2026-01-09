import fitz
from pptx import Presentation

def load_pdf(path: str) -> str:
    text = ""
    doc = fitz.open(path)
    for page in doc:
        text += page.get_text() + "\n"
    return text.strip()

def load_ppt(path: str) -> str:
    text = ""
    prs = Presentation(path)
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text.strip()
